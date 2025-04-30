import streamlit as st
import os
from PIL import Image
import base64
import requests
import json
import io
from datetime import datetime
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas
from reportlab.lib import colors
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
from database import create_user, verify_user_by_email, get_user_profile, update_user_profile

# Set page config
st.set_page_config(
    page_title="Medical Report Generator",
    page_icon="🏥",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# Initialize session state for authentication and language
if 'authenticated' not in st.session_state:
    st.session_state.authenticated = False
if 'user_email' not in st.session_state:
    st.session_state.user_email = None
if 'user_info' not in st.session_state:
    st.session_state.user_info = {}
if 'language' not in st.session_state:
    st.session_state.language = 'English'
if 'generated_report' not in st.session_state:
    st.session_state.generated_report = None

# Available languages
LANGUAGES = {
    'English': 'en',
    'Français': 'fr',
    'Español': 'es',
    'Deutsch': 'de',
    'Italiano': 'it',
    'Português': 'pt',
    'Nederlands': 'nl',
    'Polski': 'pl',
    'Русский': 'ru',
    '日本語': 'ja',
    '中文': 'zh',
    '한국어': 'ko',
    'العربية': 'ar'
}

# Language-specific prompts
LANGUAGE_PROMPTS = {
    'en': "Analyze this medical image in detail. Provide a comprehensive medical report including any visible conditions, abnormalities, or relevant medical observations. Format the response in a clear, professional manner suitable for medical documentation.",
    'fr': "Analysez cette image médicale en détail. Fournissez un rapport médical complet incluant toutes les conditions visibles, anomalies ou observations médicales pertinentes. Formatez la réponse de manière claire et professionnelle, adaptée à la documentation médicale.",
    'es': "Analice esta imagen médica en detalle. Proporcione un informe médico completo que incluya cualquier condición visible, anomalías u observaciones médicas relevantes. Formatee la respuesta de manera clara y profesional, adecuada para la documentación médica.",
    'de': "Analysieren Sie dieses medizinische Bild im Detail. Erstellen Sie einen umfassenden medizinischen Bericht, der alle sichtbaren Erkrankungen, Abnormalitäten oder relevanten medizinischen Beobachtungen enthält. Formatieren Sie die Antwort klar und professionell, geeignet für die medizinische Dokumentation.",
    'it': "Analizzare questa immagine medica in dettaglio. Fornire un rapporto medico completo che includa eventuali condizioni visibili, anomalie o osservazioni mediche rilevanti. Formattare la risposta in modo chiaro e professionale, adatto alla documentazione medica.",
    'pt': "Analise esta imagem médica em detalhes. Forneça um relatório médico abrangente incluindo quaisquer condições visíveis, anomalias ou observações médicas relevantes. Formate a resposta de maneira clara e profissional, adequada para documentação médica.",
    'nl': "Analyseer dit medische beeld in detail. Geef een uitgebreid medisch rapport met alle zichtbare aandoeningen, afwijkingen of relevante medische observaties. Formatteer het antwoord op een duidelijke, professionele manier die geschikt is voor medische documentatie.",
    'pl': "Przeanalizuj szczegółowo ten obraz medyczny. Przedstaw kompleksowy raport medyczny zawierający wszystkie widoczne schorzenia, nieprawidłowości lub istotne obserwacje medyczne. Sformatuj odpowiedź w jasny, profesjonalny sposób odpowiedni do dokumentacji medycznej.",
    'ru': "Проанализируйте это медицинское изображение подробно. Предоставьте комплексный медицинский отчет, включающий все видимые состояния, аномалии или соответствующие медицинские наблюдения. Отформатируйте ответ четко и профессионально, подходяще для медицинской документации.",
    'ja': "この医療画像を詳細に分析してください。見られる症状、異常、または関連する医療観察を含む包括的な医療報告書を提供してください。医療文書に適した明確で専門的な形式で回答を構成してください。",
    'zh': "详细分析这张医疗图像。提供一份全面的医疗报告，包括任何可见的病症、异常或相关医疗观察。以清晰、专业的方式格式化回答，适合医疗文档。",
    'ko': "이 의료 이미지를 자세히 분석하세요. 보이는 모든 상태, 이상 또는 관련 의료 관찰을 포함하는 포괄적인 의료 보고서를 제공하세요. 의료 문서에 적합한 명확하고 전문적인 방식으로 응답을 구성하세요.",
    'ar': "قم بتحليل هذه الصورة الطبية بالتفصيل. قدم تقريراً طبياً شاملاً يتضمن أي حالات مرئية أو تشوهات أو ملاحظات طبية ذات صلة. قم بتنسيق الرد بطريقة واضحة ومهنية مناسبة للتوثيق الطبي."
}

# List of potential medical indications (for auto-generation)
MEDICAL_INDICATIONS = {
    'en': [
        "Follow-up examination after treatment",
        "Suspected fracture",
        "Routine check-up",
        "Persistent pain in affected area",
        "Post-surgical evaluation",
        "Monitoring disease progression",
        "Evaluation of treatment efficacy",
        "Assessment of abnormality detected in previous exam",
        "Pre-operative assessment",
        "Investigation of unexplained symptoms"
    ],
    'fr': [
        "Examen de suivi après traitement",
        "Suspicion de fracture",
        "Contrôle de routine",
        "Douleur persistante dans la zone affectée",
        "Évaluation post-chirurgicale",
        "Suivi de l'évolution de la maladie",
        "Évaluation de l'efficacité du traitement",
        "Évaluation d'une anomalie détectée lors d'un examen précédent",
        "Évaluation préopératoire",
        "Investigation de symptômes inexpliqués"
    ]
}

# Enhanced Custom CSS
st.markdown("""
    <style>
    /* Main app styling */
    .main {
        background-color: #f8f9fa;
        padding: 2rem;
        font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
    }
    
    /* Header styling */
    h1, h2, h3 {
        color: #2c3e50;
        font-weight: 600;
    }
    
    h1 {
        font-size: 2.5rem;
        margin-bottom: 1.5rem;
        background: linear-gradient(90deg, #4b6cb7 0%, #182848 100%);
        -webkit-background-clip: text;
        -webkit-text-fill-color: transparent;
        padding: 10px 0;
    }
    
    /* Button styling */
    .stButton>button {
        width: 100%;
        height: 3em;
        background: linear-gradient(90deg, #36D1DC 0%, #5B86E5 100%);
        color: white;
        border: none;
        border-radius: 8px;
        font-weight: 500;
        transition: all 0.3s ease;
        box-shadow: 0 4px 6px rgba(50, 50, 93, 0.11), 0 1px 3px rgba(0, 0, 0, 0.08);
    }
    
    .stButton>button:hover {
        transform: translateY(-2px);
        box-shadow: 0 7px 14px rgba(50, 50, 93, 0.1), 0 3px 6px rgba(0, 0, 0, 0.08);
        background: linear-gradient(90deg, #5B86E5 0%, #36D1DC 100%);
    }
    
    /* Form styling */
    .stTextInput>div>div>input, .stSelectbox>div>div>div {
        border-radius: 6px;
        border: 1px solid #e0e0e0;
        padding: 12px;
        box-shadow: none !important;
        font-size: 16px;
    }
    
    .stTextInput>div>div>input:focus, .stSelectbox>div>div>div:focus {
        border: 1px solid #5B86E5;
        box-shadow: 0 0 0 1px #5B86E5 !important;
    }
    
    /* Authentication container */
    .auth-container {
        max-width: 450px;
        margin: 2rem auto;
        padding: 2rem;
        background-color: white;
        border-radius: 12px;
        box-shadow: 0 10px 25px rgba(0,0,0,0.05);
        border: 1px solid #f0f0f0;
    }
    
    /* Tabs styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
        margin-bottom: 0.5rem;
    }
    
    .stTabs [data-baseweb="tab"] {
        padding: 10px 24px;
        background-color: #f8f9fa;
        border-radius: 8px 8px 0 0;
        border: none;
        color: #4a4a4a;
        font-weight: 500;
    }
    
    .stTabs [aria-selected="true"] {
        background-color: #5B86E5 !important;
        color: white !important;
    }
    
    /* File uploader */
    .uploadedFile {
        border: 1px dashed #5B86E5;
        border-radius: 10px;
        padding: 20px;
        text-align: center;
        background-color: #f7fafc;
    }
    
    /* Card styling for report */
    .report-card {
        background-color: white;
        padding: 2rem;
        border-radius: 12px;
        box-shadow: 0 5px 15px rgba(0,0,0,0.05);
        margin: 1.5rem 0;
        border-left: 5px solid #5B86E5;
    }
    
    /* Success/error messages */
    .success-message {
        padding: 1rem;
        background-color: #d4edda;
        color: #155724;
        border-radius: 6px;
        margin: 1rem 0;
    }
    
    .error-message {
        padding: 1rem;
        background-color: #f8d7da;
        color: #721c24;
        border-radius: 6px;
        margin: 1rem 0;
    }
    
    /* Language selector */
    .language-selector {
        background-color: white;
        padding: 1rem;
        border-radius: 8px;
        box-shadow: 0 2px 5px rgba(0,0,0,0.05);
        margin-bottom: 1rem;
    }
    
    /* Footer */
    .footer {
        margin-top: 3rem;
        padding-top: 1.5rem;
        border-top: 1px solid #e0e0e0;
        text-align: center;
        color: #6c757d;
    }
    
    /* Download button */
    .download-btn {
        background: linear-gradient(90deg, #11998e 0%, #38ef7d 100%);
        margin-top: 1rem;
    }
    
    /* Download options */
    .download-options {
        display: flex;
        justify-content: center;
        gap: 10px;
        margin-top: 20px;
    }
    
    .format-button {
        padding: 8px 16px;
        background-color: #5B86E5;
        color: white;
        border-radius: 4px;
        text-decoration: none;
        font-size: 14px;
        font-weight: 500;
        transition: all 0.2s ease;
    }
    
    .format-button:hover {
        background-color: #4a6bba;
        transform: translateY(-2px);
    }
    
    /* Profile section */
    .profile-section {
        display: flex;
        justify-content: space-between;
        align-items: center;
        background-color: white;
        padding: 1rem;
        border-radius: 8px;
        box-shadow: 0 2px 5px rgba(0,0,0,0.05);
        margin-bottom: 1.5rem;
    }
    
    /* Profile form */
    .profile-form {
        background-color: white;
        padding: 20px;
        border-radius: 8px;
        box-shadow: 0 2px 10px rgba(0,0,0,0.05);
        margin-bottom: 20px;
    }
    </style>
""", unsafe_allow_html=True)

def login(username, password):
    user = verify_user_by_email(username, password)
    if user:
        st.session_state.authenticated = True
        st.session_state.username = username
        st.session_state.display_name = get_user_profile(username).get('username', username)
        return True
    return False

def logout():
    st.session_state.authenticated = False
    st.session_state.user_email = None
    st.session_state.user_info = {}
    st.session_state.generated_report = None

# Function to generate PDF report
def generate_pdf(analysis, patient_info, image_bytes=None):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    
    # Create custom styles
    title_style = ParagraphStyle(
        'Title',
        parent=styles['Heading1'],
        fontSize=16,
        textColor=colors.darkblue,
        spaceAfter=12
    )
    
    header_style = ParagraphStyle(
        'Header',
        parent=styles['Heading2'],
        fontSize=14,
        textColor=colors.darkblue,
        spaceAfter=10
    )
    
    normal_style = ParagraphStyle(
        'Normal',
        parent=styles['Normal'],
        fontSize=11,
        spaceAfter=8
    )
    
    # Create the story
    story = []
    
    # Title
    title = Paragraph("RAPPORT D'ANALYSE MÉDICALE", title_style)
    story.append(title)
    story.append(Spacer(1, 12))
    
    # Date
    date_str = datetime.now().strftime("%d/%m/%Y à %H:%M")
    date_paragraph = Paragraph(f"Date du rapport: {date_str}", normal_style)
    story.append(date_paragraph)
    story.append(Spacer(1, 12))
    
    # Patient info
    story.append(Paragraph("INFORMATIONS PATIENT", header_style))
    story.append(Paragraph(f"Nom du patient: {patient_info['name']}", normal_style))
    story.append(Paragraph(f"Numéro d'identification: {patient_info['id']}", normal_style))
    story.append(Paragraph(f"Indications: {patient_info['indications']}", normal_style))
    story.append(Spacer(1, 12))
    
    # Add image if available
    if image_bytes:
        img_stream = io.BytesIO(image_bytes)
        img = RLImage(img_stream, width=300, height=200)
        story.append(Paragraph("IMAGE MÉDICALE", header_style))
        story.append(img)
        story.append(Spacer(1, 12))
    
    # Analysis
    story.append(Paragraph("ANALYSE MÉDICALE", header_style))
    analysis_paragraphs = analysis.split('\n')
    for para in analysis_paragraphs:
        if para.strip():
            story.append(Paragraph(para, normal_style))
            story.append(Spacer(1, 6))
    
    # Disclaimer
    story.append(Spacer(1, 20))
    disclaimer = Paragraph("Ce rapport est généré par une intelligence artificielle et doit être examiné par un professionnel de la santé qualifié. Il ne constitue pas un avis médical professionnel.", 
                          ParagraphStyle('Disclaimer', parent=normal_style, fontSize=9, textColor=colors.grey))
    story.append(disclaimer)
    
    # Build PDF
    doc.build(story)
    buffer.seek(0)
    return buffer

# Function to generate DOCX report
def generate_docx(analysis, patient_info):
    doc = Document()
    
    # Title
    doc.add_heading('RAPPORT D\'ANALYSE MÉDICALE', 0)
    
    # Date
    date_str = datetime.now().strftime("%d/%m/%Y à %H:%M")
    doc.add_paragraph(f"Date du rapport: {date_str}")
    
    # Patient info
    doc.add_heading('INFORMATIONS PATIENT', level=1)
    doc.add_paragraph(f"Nom du patient: {patient_info['name']}")
    doc.add_paragraph(f"Numéro d'identification: {patient_info['id']}")
    doc.add_paragraph(f"Indications: {patient_info['indications']}")
    
    # Analysis
    doc.add_heading('ANALYSE MÉDICALE', level=1)
    for para in analysis.split('\n'):
        if para.strip():
            doc.add_paragraph(para)
    
    # Disclaimer
    doc.add_paragraph('Ce rapport est généré par une intelligence artificielle et doit être examiné par un professionnel de la santé qualifié. Il ne constitue pas un avis médical professionnel.').italic = True
    
    # Save to BytesIO
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

# Function to generate PPTX report
def generate_pptx(analysis, patient_info, image_bytes=None):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "RAPPORT D'ANALYSE MÉDICALE"
    date_str = datetime.now().strftime("%d/%m/%Y")
    subtitle.text = f"Généré le {date_str}"
    
    # Patient info slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Informations Patient"
    
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    p = tf.paragraphs[0]
    p.text = f"Nom du patient: {patient_info['name']}"
    
    p = tf.add_paragraph()
    p.text = f"Numéro d'identification: {patient_info['id']}"
    
    p = tf.add_paragraph()
    p.text = f"Indications: {patient_info['indications']}"
    
    # Image slide if available
    if image_bytes:
        img_slide_layout = prs.slide_layouts[6]  # Blank slide
        slide = prs.slides.add_slide(img_slide_layout)
        
        # Add title manually
        left = top = Inches(1)
        width = Inches(8)
        height = Inches(1)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        p = tf.add_paragraph()
        p.text = "Image Médicale"
        p.font.size = Pt(24)
        p.font.bold = True
        
        # Save image to temp file for adding to PowerPoint
        img_stream = io.BytesIO(image_bytes)
        img_path = "temp_image.png"
        with open(img_path, "wb") as f:
            f.write(img_stream.getvalue())
        
        # Add image to slide
        left = Inches(2)
        top = Inches(2)
        pic = slide.shapes.add_picture(img_path, left, top, height=Inches(4))
        
        # Remove temp file
        os.remove(img_path)
    
    # Analysis slide(s)
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Analyse Médicale"
    
    # Break analysis into chunks if it's long
    analysis_paragraphs = analysis.split('\n')
    chunks = [analysis_paragraphs[i:i+5] for i in range(0, len(analysis_paragraphs), 5)]
    
    # First chunk on the first slide
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    for i, para in enumerate(chunks[0]):
        if para.strip():
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = para
    
    # Additional slides for remaining chunks
    for chunk in chunks[1:]:
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = "Analyse Médicale (Suite)"
        
        shapes = slide.shapes
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for i, para in enumerate(chunk):
            if para.strip():
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = para
    
    # Disclaimer slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Avertissement"
    
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.paragraphs[0]
    p.text = "Ce rapport est généré par une intelligence artificielle et doit être examiné par un professionnel de la santé qualifié. Il ne constitue pas un avis médical professionnel."
    
    # Save to BytesIO
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# Generate random medical indications based on language
def generate_medical_indications(language_code):
    lang = language_code if language_code in MEDICAL_INDICATIONS else 'en'
    import random
    num_indications = random.randint(1, 3)
    selected_indications = random.sample(MEDICAL_INDICATIONS[lang], num_indications)
    return ", ".join(selected_indications)

# Authentication UI
if not st.session_state.authenticated:
    # Create a centered container with a medical logo
    col1, col2, col3 = st.columns([1,2,1])
    with col2:
        st.markdown("""
            <div style='text-align: center; padding: 20px;'>
                <h1>🏥 Medical Report Generator</h1>
                <p style='font-size: 18px; color: #6c757d; margin-bottom: 30px;'>
                    Professional AI-powered analysis for medical imaging
                </p>
            </div>
        """, unsafe_allow_html=True)
    
    st.markdown("<div class='auth-container'>", unsafe_allow_html=True)
    
    # Create tabs for Login and Sign Up
    tab1, tab2 = st.tabs(["Login", "Sign Up"])
    
    with tab1:
        st.markdown("<h3 style='text-align: center; margin-bottom: 20px;'>Welcome Back</h3>", unsafe_allow_html=True)
        
        with st.form("login_form"):
            email = st.text_input("Email")
            password = st.text_input("Password", type="password")
            st.markdown("<br>", unsafe_allow_html=True)
            submit = st.form_submit_button("Login")
            
            if submit:
                if login(email, password):
                    st.markdown("""
                        <div class='success-message'>
                            Login successful! Redirecting...
                        </div>
                    """, unsafe_allow_html=True)
                    st.rerun()
                else:
                    st.markdown("""
                        <div class='error-message'>
                            Invalid email or password
                        </div>
                    """, unsafe_allow_html=True)
    
    with tab2:
        st.markdown("<h3 style='text-align: center; margin-bottom: 20px;'>Create New Account</h3>", unsafe_allow_html=True)
        
        with st.form("signup_form"):
            new_username = st.text_input("Username")
            new_email = st.text_input("Email")
            new_password = st.text_input("Choose Password", type="password")
            confirm_password = st.text_input("Confirm Password", type="password")
            st.markdown("<br>", unsafe_allow_html=True)
            submit_signup = st.form_submit_button("Sign Up")
            
            if submit_signup:
                if new_password != confirm_password:
                    st.markdown("""
                        <div class='error-message'>
                            Passwords do not match!
                        </div>
                    """, unsafe_allow_html=True)
                elif len(new_password) < 6:
                    st.markdown("""
                        <div class='error-message'>
                            Password must be at least 6 characters long!
                        </div>
                    """, unsafe_allow_html=True)
                else:
                    if create_user(new_username, new_password, new_email):
                        st.markdown("""
                            <div class='success-message'>
                                Account created successfully! Please login.
                            </div>
                        """, unsafe_allow_html=True)
                    else:
                        st.markdown("""
                            <div class='error-message'>
                                Username or email already exists!
                            </div>
                        """, unsafe_allow_html=True)
    
    st.markdown("</div>", unsafe_allow_html=True)
    
    # Add some testimonials or features section
    st.markdown("""
        <div style='max-width: 800px; margin: 3rem auto; text-align: center;'>
            <h3>Why Choose Our Medical Report Generator?</h3>
            <div style='display: flex; justify-content: space-between; margin-top: 30px;'>
                <div style='flex: 1; padding: 20px; background-color: white; border-radius: 8px; margin: 0 10px; box-shadow: 0 5px 15px rgba(0,0,0,0.05);'>
                    <h4>🔍 Advanced AI Analysis</h4>
                    <p>Powered by Google's Gemini AI for accurate medical image interpretation</p>
                </div>
                <div style='flex: 1; padding: 20px; background-color: white; border-radius: 8px; margin: 0 10px; box-shadow: 0 5px 15px rgba(0,0,0,0.05);'>
                    <h4>🌍 Multilingual Support</h4>
                    <p>Generate reports in 13 different languages for global accessibility</p>
                </div>
                <div style='flex: 1; padding: 20px; background-color: white; border-radius: 8px; margin: 0 10px; box-shadow: 0 5px 15px rgba(0,0,0,0.05);'>
                    <h4>🔒 Secure & Private</h4>
                    <p>Your medical data is processed with the highest security standards</p>
                </div>
            </div>
        </div>
    """, unsafe_allow_html=True)
    
    st.stop()

# Main application (only shown after authentication)
st.markdown("""
    <h1>🏥 Medical Report Generator</h1>
""", unsafe_allow_html=True)

# User profile section with logout button implemented in Python (not JavaScript)
col1, col2 = st.columns([3, 1])
with col1:
    st.markdown(f"""
        <div style='background-color: white; padding: 15px; border-radius: 8px; box-shadow: 0 2px 5px rgba(0,0,0,0.05);'>
            <h3 style='margin-bottom: 5px;'>Welcome, {st.session_state.display_name}!</h3>
            <p style='color: #6c757d;'>Let's generate some professional medical reports today.</p>
        </div>
    """, unsafe_allow_html=True)
    
with col2:
    # Implemented logout correctly with Python
    if st.button('Logout', key='logout_button'):
        logout()
        st.rerun()

# Create tabs for Profile and Report Generator
tabs = st.tabs(["Profile", "Report Generator"])

with tabs[0]:
    st.markdown("<h3>Patient Profile</h3>", unsafe_allow_html=True)
    
    # Display existing user info and allow editing
    with st.form("profile_form", clear_on_submit=False):
        st.markdown("<div class='profile-form'>", unsafe_allow_html=True)
        
        col1, col2 = st.columns(2)
        with col1:
            first_name = st.text_input("First Name", value=st.session_state.user_info.get('first_name', ''))
            last_name = st.text_input("Last Name", value=st.session_state.user_info.get('last_name', ''))
            dob = st.date_input("Date of Birth", value=None)
            
        with col2:
            patient_id = st.text_input("Patient ID", value=st.session_state.user_info.get('patient_id', ''))
            gender = st.selectbox("Gender", ["Male", "Female", "Other", "Prefer not to say"], index=0)
            phone = st.text_input("Phone Number", value=st.session_state.user_info.get('phone', ''))
        
        medical_history = st.text_area("Medical History", value=st.session_state.user_info.get('medical_history', ''))
        current_medications = st.text_area("Current Medications", value=st.session_state.user_info.get('medications', ''))
        
        submit_profile = st.form_submit_button("Update Profile")
        
        if submit_profile:
            # Update profile in session state and database
            profile_data = {
                'first_name': first_name,
                'last_name': last_name,
                'dob': str(dob),
                'patient_id': patient_id,
                'gender': gender,
                'phone': phone,
                'medical_history': medical_history,
                'medications': current_medications
            }
            update_user_profile(st.session_state.user_email, profile_data)
            st.session_state.user_info.update(profile_data)
            st.success("Profile updated successfully!")
            st.rerun()
        
        st.markdown("</div>", unsafe_allow_html=True)

with tabs[1]:
    # Two column layout for better organization
    col1, col2 = st.columns([1, 2])

    with col1:
        st.markdown("<div class='language-selector'>", unsafe_allow_html=True)
        st.markdown("<h4>Report Settings</h4>", unsafe_allow_html=True)
        
        # Language selection
        selected_language = st.selectbox(
            "Select Report Language",
            options=list(LANGUAGES.keys()),
            index=list(LANGUAGES.keys()).index(st.session_state.language)
        )
        st.session_state.language = selected_language
        
        # Patient info for report
        st.markdown("<h4>Patient Information</h4>", unsafe_allow_html=True)
        
        full_name = f"{st.session_state.user_info.get('first_name', '')} {st.session_state.user_info.get('last_name', '')}"
        patient_name = st.text_input("Patient Name", value=full_name)
        patient_id_value = st.text_input("Patient ID", value=st.session_state.user_info.get('patient_id', ''))
        
        # API Key input
        st.markdown("<h4>API Settings</h4>", unsafe_allow_html=True)
        api_key = st.text_input("Enter your Google AI Studio API Key:", type="password")
        
        # Generate random indications or allow custom input
        st.markdown("<h4>Medical Indications</h4>", unsafe_allow_html=True)
        use_auto_indications = st.checkbox("Auto-generate indications", value=True)
        
        language_code = LANGUAGES[selected_language]
        
        if use_auto_indications:
            indications = generate_medical_indications(language_code)
            st.text_area("Generated Indications", value=indications, height=100, key="auto_indications", disabled=True)
        else:
            indications = st.text_area("Enter Medical Indications", height=100, key="manual_indications")
        
        st.markdown("</div>", unsafe_allow_html=True)

    with col2:
        # File uploader with better styling
        st.markdown("<h3>Upload Medical Image</h3>", unsafe_allow_html=True)
        
        uploaded_file = st.file_uploader("Choose a medical image...", type=['jpg', 'jpeg', 'png'])
        
        if uploaded_file is not None:
            # Display the uploaded image in a nice frame
            st.markdown("<div class='uploadedFile'>", unsafe_allow_html=True)
            image = Image.open(uploaded_file)
            st.image(image, caption="Uploaded Medical Image", use_column_width=True)
            st.markdown("</div>", unsafe_allow_html=True)
            
            # Convert image to base64
            image_bytes = uploaded_file.getvalue()
            image_base64 = base64.b64encode(image_bytes).decode()
            
            # Analysis button
            if api_key:
                if st.button("Generate Report", key="generate_btn"):
                    with st.spinner("Analyzing image..."):
                        try:
                            # Prepare the request
                            url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={api_key}"
                            
                            headers = {
                                "Content-Type": "application/json"
                            }
                            
                            # Get the appropriate prompt for the selected language
                            language_code = LANGUAGES[selected_language]
                            prompt = LANGUAGE_PROMPTS[language_code]
                            
                            data = {
                                "contents": [{
                                    "parts": [
                                        {
                                            "inline_data": {
                                                "mime_type": uploaded_file.type,
                                                "data": image_base64
                                            }
                                        },
                                        {
                                            "text": prompt
                                        }
                                    ]
                                }]
                            }
                            
                            # Make the API request
                            response = requests.post(url, headers=headers, json=data)
                            response.raise_for_status()
                            
                            # Extract and display the analysis
                            result = response.json()
                            analysis = result['candidates'][0]['content']['parts'][0]['text']
                            
                            # Save report info in session state
                            patient_info = {
                                'name': patient_name,
                                'id': patient_id_value,
                                'indications': indications
                            }
                            
                            # Store in session state
                            st.session_state.generated_report = {
                                'analysis': analysis,
                                'patient_info': patient_info,
                                'image_bytes': image_bytes,
                                'language': selected_language
                            }
                            
                            # Display the report in a nice format
                            st.markdown("<div class='report-card'>", unsafe_allow_html=True)
                            st.markdown(f"<h3>📋 Medical Analysis Report ({selected_language})</h3>", unsafe_allow_html=True)
                            st.markdown("<hr style='margin-bottom: 20px;'>", unsafe_allow_html=True)
                            
                            # Add patient info to displayed report
                            st.markdown(f"<p><strong>Patient:</strong> {patient_name}</p>", unsafe_allow_html=True)
                            st.markdown(f"<p><strong>Patient ID:</strong> {patient_id_value}</p>", unsafe_allow_html=True)
                            st.markdown(f"<p><strong>Indications:</strong> {indications}</p>", unsafe_allow_html=True)
                            st.markdown("<hr style='margin-bottom: 20px;'>", unsafe_allow_html=True)
                            
                            # Format the analysis into a single coherent block
                            formatted_analysis = analysis.replace('**', '')  # Remove asterisks
                            formatted_analysis = formatted_analysis.replace('📝 Medical Analysis', '')  # Remove repeated headers
                            formatted_analysis = formatted_analysis.replace('📊 Analysis Accuracy', '')  # Remove repeated headers
                            formatted_analysis = formatted_analysis.replace('💊 Recommended Medications', '')  # Remove repeated headers
                            formatted_analysis = formatted_analysis.replace('⚠️ Severity Assessment', '')  # Remove repeated headers
                            
                            # Display analysis in a clean, continuous format
                            st.markdown(f"""
                                <div style='font-family: Arial, sans-serif; line-height: 1.6; font-size: 16px; color: #333;'>
                                    {formatted_analysis}
                                </div>
                            """, unsafe_allow_html=True)
                            
                            st.markdown("</div>", unsafe_allow_html=True)
                            
                            # Add download options for the report with better styling
                            st.markdown("<h4>Download Report</h4>", unsafe_allow_html=True)
                            col1, col2, col3, col4 = st.columns(4)
                            with col1:
                                pdf_buffer = generate_pdf(analysis, patient_info, image_bytes)
                                st.download_button(
                                    label="PDF Format",
                                    data=pdf_buffer,
                                    file_name=f"medical_report_{language_code}.pdf",
                                    mime="application/pdf",
                                    key="pdf_btn"
                                )
                            with col2:
                                docx_buffer = generate_docx(analysis, patient_info)
                                st.download_button(
                                    label="Word Format",
                                    data=docx_buffer,
                                    file_name=f"medical_report_{language_code}.docx",
                                    mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                                    key="docx_btn"
                                )
                            with col3:
                                pptx_buffer = generate_pptx(analysis, patient_info, image_bytes)
                                st.download_button(
                                    label="PowerPoint",
                                    data=pptx_buffer,
                                    file_name=f"medical_report_{language_code}.pptx",
                                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                                    key="pptx_btn"
                                )
                            with col4:
                                report_text = f"""MEDICAL ANALYSIS REPORT
                                
Patient: {patient_name}
Patient ID: {patient_id_value}
Indications: {indications}

ANALYSIS:
{analysis}

Generated on {datetime.now().strftime("%d/%m/%Y at %H:%M")}
This report is generated by AI and should be reviewed by a qualified healthcare professional.
"""
                                st.download_button(
                                    label="Text Format",
                                    data=report_text,
                                    file_name=f"medical_report_{language_code}.txt",
                                    mime="text/plain",
                                    key="txt_btn"
                                )
                            
                        except Exception as e:
                            st.markdown(f"""
                                <div class='error-message'>
                                    <h4>⚠️ Error</h4>
                                    <p>{str(e)}</p>
                                    <p>Please check your API key and try again.</p>
                                </div>
                            """, unsafe_allow_html=True)
            elif uploaded_file is not None:
                st.warning("Please enter your Google AI Studio API key to proceed with the analysis.")
        else:
            # Display a placeholder or instructions
            st.markdown("""
                <div style='background-color: #f8f9fa; border: 2px dashed #dddddd; border-radius: 10px; padding: 50px 20px; text-align: center; margin-top: 20px;'>
                    <img src="https://via.placeholder.com/150/5B86E5/FFFFFF?text=📊" style='width: 100px; height: 100px; opacity: 0.7;'>
                    <h3 style='margin-top: 20px; color: #6c757d;'>Upload a Medical Image</h3>
                    <p style='color: #6c757d;'>Supported formats: JPG, JPEG, PNG</p>
                </div>
            """, unsafe_allow_html=True)

# Footer with better styling
st.markdown("""
    <div class='footer'>
        <div style='display: flex; justify-content: center; margin-bottom: 15px;'>
            <div style='margin: 0 15px;'>Terms of Service</div>
            <div style='margin: 0 15px;'>Privacy Policy</div>
            <div style='margin: 0 15px;'>FAQs</div>
            <div style='margin: 0 15px;'>Contact</div>
        </div>
        <p>© 2025 Medical Report Generator. Built with ❤️ using Streamlit and Google's Gemini AI</p>
        <p style='font-size: 12px; color: #adb5bd; margin-top: 10px;'>This tool is intended for informational purposes only and should not replace professional medical advice.</p>
    </div>
""", unsafe_allow_html=True)