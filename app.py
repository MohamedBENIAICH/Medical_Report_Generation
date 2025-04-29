import streamlit as st
import os
from PIL import Image
import base64
import requests
import json
from database import create_user, verify_user_by_email as verify_user
import io
from datetime import datetime
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib import colors
from docx import Document
from pptx import Presentation
from pptx.util import Inches, Pt
import uuid
from reportlab.lib.pagesizes import letter
from reportlab.lib import colors
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Image as RLImage
from reportlab.pdfgen.canvas import Canvas
import qrcode
from google_auth import handle_google_signin
import time

# Set page config
st.set_page_config(
    page_title="Medical Report Generator",
    page_icon="🏥",
    layout="wide",
    initial_sidebar_state="collapsed"
)

# Initialize session state for authentication and language
if 'authenticated' not in st.session_state:
    st.session_state.authenticated = False
if 'username' not in st.session_state:
    st.session_state.username = None
if 'language' not in st.session_state:
    st.session_state.language = 'English'
if 'analysis_result' not in st.session_state:
    st.session_state.analysis_result = None
if 'patient_info' not in st.session_state:
    st.session_state.patient_info = None
if 'image_bytes' not in st.session_state:
    st.session_state.image_bytes = None

# Available languages
LANGUAGES = {
    'English': 'en',
    'Français': 'fr',
    'Español': 'es',
    'Deutsch': 'de',
    'Italiano': 'it',
    'Português': 'pt',
    'Nederlands': 'nl',
    'Polski': 'pl',
    'Русский': 'ru',
    '日本語': 'ja',
    '中文': 'zh',
    '한국어': 'ko',
    'العربية': 'ar'
}

# Language-specific prompts
LANGUAGE_PROMPTS = {
    'en': "Analyze this medical image in detail. Provide a comprehensive medical report including any visible conditions, abnormalities, or relevant medical observations. Format the response in a clear, professional manner suitable for medical documentation.",
    'fr': "Analysez cette image médicale en détail. Fournissez un rapport médical complet incluant toutes les conditions visibles, anomalies ou observations médicales pertinentes. Formatez la réponse de manière claire et professionnelle, adaptée à la documentation médicale.",
    'es': "Analice esta imagen médica en detalle. Proporcione un informe médico completo que incluya cualquier condición visible, anomalías u observaciones médicas relevantes. Formatee la respuesta de manera clara et professionnelle, adecuada para la documentación médica.",
    'de': "Analysieren Sie dieses medizinische Bild im Detail. Erstellen Sie einen umfassenden medizinischen Bericht, der alle sichtbaren Erkrankungen, Abnormalitäten oder relevanten medizinischen Beobachtungen enthält. Formatieren Sie die Antwort klar und professionell, geeignet für die medizinische Dokumentation.",
    'it': "Analizzare questa immagine medica in dettaglio. Fornire un rapporto medico completo che includa eventuali condizioni visibili, anomalie o osservazioni mediche rilevanti. Formattare la risposta in modo chiaro e professionale, adatto alla documentazione medica.",
    'pt': "Analise esta imagem médica em detalhes. Forneça um relatório médico abrangente incluindo quaisquer condições visíveis, anomalias ou observações médicas relevantes. Formate a resposta de maneira clara e profissional, adequada para documentação médica.",
    'nl': "Analyseer dit medische beeld in detail. Geef een uitgebreid medisch rapport met alle zichtbare aandoeningen, afwijkingen of relevante medische observaties. Formatteer het antwoord op een duidelijke, professionele manier die geschikt is voor medische documentatie.",
    'pl': "Przeanalizuj szczegółowo ten obraz medyczny. Przedstaw kompleksowy raport medyczny zawierający wszystkie widoczne schorzenia, nieprawidłowości lub istotne obserwacje medyczne. Sformatuj odpowiedź w jasny, profesjonalny sposób odpowiedni do dokumentacji medycznej.",
    'ru': "Проанализируйте это медицинское изображение подробно. Предоставьте комплексный медицинский отчет, включающий все видимые состояния, аномалии или соответствующие медицинские наблюдения. Отформатируйте ответ четко и профессионально, подходяще для медицинской документации.",
    'ja': "この医療画像を詳細に分析してください。見られる症状、異常、または関連する医療観察を含む包括的な医療報告書を提供してください。医療文書に適した明確で専門的な形式で回答を構成してください。",
    'zh': "详细分析这张医疗图像。提供一份全面的医疗报告，包括任何可见的病症、异常或相关医疗观察。以清晰、专业的方式格式化回答，适合医疗文档。",
    'ko': "이 의료 이미지를 자세히 분석하세요. 보이는 모든 상태, 이상 또는 관련 의료 관찰을 포함하는 포괄적인 의료 보고서를 제공하세요. 의료 문서에 적합한 명확하고 전문적인 방식으로 응답을 구성하세요。",
    'ar': "قم بتحليل هذه الصورة الطبية بالتفصيل. قدم تقريراً طبياً شاملاً يتضمن أي حالات مرئية أو تشوهات أو ملاحظات طبية ذات صلة. قم بتنسيق الرد بطريقة واضحة ومهنية مناسبة للتوثيق الطبي."
}

# Enhanced Custom CSS
st.markdown("""
    <style>
    /* Main app styling */
    .main {
        background-color: #f8fafc;
        padding: 2rem;
        font-family: 'Inter', 'Segoe UI', sans-serif;
    }
    
    /* Header styling */
    h1, h2, h3, h4 {
        color: #1e3a8a;
        font-weight: 600;
        font-family: 'Inter', sans-serif;
    }
    
    h1 {
        font-size: 2.2rem;
        margin-bottom: 1.5rem;
        color: #1e3a8a;
        padding: 10px 0;
        border-bottom: 2px solid #e2e8f0;
    }
    
    /* Button styling */
    .stButton>button {
        width: 100%;
        height: 2.8em;
        background-color: #2563eb;
        color: white;
        border: none;
        border-radius: 6px;
        font-weight: 500;
        transition: all 0.2s ease;
        box-shadow: 0 2px 4px rgba(37, 99, 235, 0.1);
    }
    
    .stButton>button:hover {
        background-color: #1d4ed8;
        box-shadow: 0 4px 6px rgba(37, 99, 235, 0.2);
    }
    
    /* Form styling */
    .stTextInput>div>div>input, .stSelectbox>div>div>div, .stTextArea>div>div>textarea {
        border-radius: 6px;
        border: 1px solid #e2e8f0;
        padding: 12px;
        box-shadow: none !important;
        font-size: 14px;
        background-color: white;
    }
    
    .stTextInput>div>div>input:focus, .stSelectbox>div>div>div:focus, .stTextArea>div>div>textarea:focus {
        border: 1px solid #2563eb;
        box-shadow: 0 0 0 2px rgba(37, 99, 235, 0.1) !important;
    }
    
    /* Authentication container */
    .auth-container {
        max-width: 450px;
        margin: 2rem auto;
        padding: 2rem;
        background-color: white;
        border-radius: 8px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
        border: 1px solid #e2e8f0;
    }
    
    /* Tabs styling */
    .stTabs [data-baseweb="tab-list"] {
        gap: 8px;
        margin-bottom: 1rem;
    }
    
    .stTabs [data-baseweb="tab"] {
        padding: 10px 20px;
        background-color: #f8fafc;
        border-radius: 6px;
        border: 1px solid #e2e8f0;
        color: #4b5563;
        font-weight: 500;
        font-size: 14px;
    }
    
    .stTabs [aria-selected="true"] {
        background-color: #2563eb !important;
        color: white !important;
        border-color: #2563eb !important;
    }
    
    /* File uploader */
    .uploadedFile {
        border: 2px dashed #2563eb;
        border-radius: 8px;
        padding: 20px;
        text-align: center;
        background-color: #f8fafc;
    }
    
    /* Card styling for report */
    .report-card {
        background-color: white;
        padding: 2rem;
        border-radius: 8px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
        margin: 1.5rem 0;
        border-left: 4px solid #2563eb;
    }
    
    /* Success/error messages */
    .success-message {
        padding: 1rem;
        background-color: #dcfce7;
        color: #166534;
        border-radius: 6px;
        margin: 1rem 0;
        border: 1px solid #bbf7d0;
    }
    
    .error-message {
        padding: 1rem;
        background-color: #fee2e2;
        color: #991b1b;
        border-radius: 6px;
        margin: 1rem 0;
        border: 1px solid #fecaca;
    }
    
    /* Language selector */
    .language-selector {
        background-color: white;
        padding: 1.5rem;
        border-radius: 8px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
        margin-bottom: 1.5rem;
        border: 1px solid #e2e8f0;
    }
    
    /* Footer */
    .footer {
        margin-top: 3rem;
        padding-top: 1.5rem;
        border-top: 1px solid #e2e8f0;
        text-align: center;
        color: #4b5563;
        font-size: 14px;
    }
    
    /* Download button */
    .download-btn {
        background-color: #2563eb;
        margin-top: 1rem;
    }
    
    /* Profile section */
    .profile-section {
        display: flex;
        justify-content: space-between;
        align-items: center;
        background-color: white;
        padding: 1.5rem;
        border-radius: 8px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
        margin-bottom: 1.5rem;
        border: 1px solid #e2e8f0;
    }
    
    /* Patient info form */
    .patient-info-form {
        background-color: white;
        padding: 1.5rem;
        border-radius: 8px;
        box-shadow: 0 4px 6px rgba(0,0,0,0.05);
        margin-bottom: 1.5rem;
        border: 1px solid #e2e8f0;
    }
    
    /* Download options */
    .download-options {
        display: flex;
        gap: 12px;
        flex-wrap: wrap;
        margin-top: 1.5rem;
    }
    
    /* Medical icons */
    .medical-icon {
        color: #2563eb;
        margin-right: 8px;
    }
    
    /* Report header */
    .report-header {
        background-color: #f8fafc;
        padding: 1rem;
        border-radius: 8px;
        margin-bottom: 1.5rem;
        border: 1px solid #e2e8f0;
    }
    
    /* Form labels */
    .stTextInput>label, .stSelectbox>label, .stTextArea>label {
        font-weight: 500;
        color: #1e3a8a;
        margin-bottom: 0.5rem;
    }
    
    /* Section headers */
    .section-header {
        color: #1e3a8a;
        font-weight: 600;
        margin-bottom: 1rem;
        padding-bottom: 0.5rem;
        border-bottom: 2px solid #e2e8f0;
    }
    </style>
""", unsafe_allow_html=True)

def login(username, password):
    user = verify_user(username, password)
    if user:
        st.session_state.authenticated = True
        st.session_state.user_email = username  # Store email
        st.session_state.username = user.get('username', username)  # Store actual username
        return True
    return False

def logout():
    st.session_state.authenticated = False
    st.session_state.user_email = None
    st.session_state.username = None
    st.session_state.analysis_result = None
    st.session_state.patient_info = None
    st.session_state.image_bytes = None

def generate_pdf(analysis, patient_info, image_bytes=None):
    buffer = io.BytesIO()
    doc = SimpleDocTemplate(buffer, pagesize=letter)
    styles = getSampleStyleSheet()
    
    # Generate a simple report ID
    report_id = f"MIR-{uuid.uuid4().hex[:8].upper()}"
    
    # Custom styles for medical report
    title_style = ParagraphStyle(
        'Title',
        parent=styles['Heading1'],
        fontSize=18,
        textColor=colors.HexColor('#1e3a8a'),
        spaceAfter=20,
        alignment=1,  # Center alignment
        fontName='Helvetica-Bold'
    )
    
    header_style = ParagraphStyle(
        'Header',
        parent=styles['Heading2'],
        fontSize=14,
        textColor=colors.HexColor('#1e3a8a'),
        spaceAfter=12,
        fontName='Helvetica-Bold'
    )
    
    subheader_style = ParagraphStyle(
        'Subheader',
        parent=styles['Heading3'],
        fontSize=12,
        textColor=colors.HexColor('#4b5563'),
        spaceAfter=8,
        fontName='Helvetica-Bold'
    )
    
    normal_style = ParagraphStyle(
        'Normal',
        parent=styles['Normal'],
        fontSize=11,
        spaceAfter=8,
        fontName='Helvetica',
        leading=14
    )
    
    # Build the PDF content
    story = []
    
    # Header with report info
    header_text = """
    <b>MEDICAL IMAGING ANALYSIS REPORT</b><br/>
    <font size="10" color="#4b5563">Generated by Medical Report Generator</font>
    """
    header = Paragraph(header_text, title_style)
    story.append(header)
    story.append(Spacer(1, 20))
    
    # Report metadata
    date_str = datetime.now().strftime("%Y-%m-%d at %H:%M")
    metadata = f"""
    <b>Report ID:</b> {report_id}<br/>
    <b>Date Generated:</b> {date_str}<br/>
    <b>Report Type:</b> Medical Imaging Analysis
    """
    story.append(Paragraph(metadata, subheader_style))
    story.append(Spacer(1, 20))
    
    # Patient Information
    story.append(Paragraph("PATIENT INFORMATION", header_style))
    patient_info_text = f"""
    <b>Name:</b> {patient_info.get('name', 'Not Provided')}<br/>
    <b>Patient ID:</b> {patient_info.get('id', 'Not Provided')}<br/>
    <b>Age:</b> {patient_info.get('age', 'Not Provided')}<br/>
    <b>Sex:</b> {patient_info.get('sex', 'Not Provided')}
    """
    story.append(Paragraph(patient_info_text, normal_style))
    story.append(Spacer(1, 20))
    
    # Image if available
    if image_bytes:
        try:
            img_stream = io.BytesIO(image_bytes)
            img = RLImage(img_stream, width=400, height=300)
            story.append(Paragraph("MEDICAL IMAGE", header_style))
            story.append(img)
            story.append(Spacer(1, 20))
        except:
            story.append(Paragraph("(Could not embed image)", normal_style))
    
    # Analysis
    story.append(Paragraph("MEDICAL ANALYSIS", header_style))
    analysis_paragraphs = analysis.split('\n')
    for para in analysis_paragraphs:
        if para.strip():
            story.append(Paragraph(para, normal_style))
            story.append(Spacer(1, 8))
    
    # Add QR code for verification
    story.append(Spacer(1, 20))
    try:
        qr = qrcode.QRCode(
            version=1,
            error_correction=qrcode.constants.ERROR_CORRECT_L,
            box_size=10,
            border=4,
        )
        qr.add_data(f"REPORT-VERIFY:{report_id}")
        qr.make(fit=True)
        
        qr_img = qr.make_image(fill_color="black", back_color="white")
        img_bytes = io.BytesIO()
        qr_img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        
        qr_image = RLImage(img_bytes, width=1.5*inch, height=1.5*inch)
        story.append(qr_image)
        story.append(Paragraph("Scan to verify report authenticity", 
                            ParagraphStyle(
                                'Caption',
                                parent=normal_style,
                                fontSize=9,
                                textColor=colors.HexColor('#6b7280'),
                                alignment=1
                            )))
    except:
        story.append(Paragraph("(Verification QR code could not be generated)", normal_style))
    
    # Footer with disclaimer
    story.append(Spacer(1, 30))
    disclaimer = Paragraph(
        "This report is generated by artificial intelligence and should be reviewed by a qualified healthcare professional. It does not constitute professional medical advice.",
        ParagraphStyle(
            'Disclaimer',
            parent=normal_style,
            fontSize=9,
            textColor=colors.HexColor('#6b7280'),
            alignment=1
        )
    )
    story.append(disclaimer)
    
    # Build the PDF
    doc.build(story)
    buffer.seek(0)
    return buffer

def generate_docx(analysis, patient_info, image_bytes=None):
    doc = Document()
    
    # Title
    doc.add_heading('MEDICAL ANALYSIS REPORT', 0)
    
    # Date
    date_str = datetime.now().strftime("%Y-%m-%d at %H:%M")
    doc.add_paragraph(f"Report Date: {date_str}")
    
    # Patient Information
    doc.add_heading('PATIENT INFORMATION', level=1)
    doc.add_paragraph(f"Patient Name: {patient_info['name']}")
    doc.add_paragraph(f"Patient ID: {patient_info['id']}")
    doc.add_paragraph(f"Age: {patient_info['age']}")
    doc.add_paragraph(f"Sex: {patient_info['sex']}")
    
    # Image if available
    if image_bytes:
        try:
            doc.add_heading('MEDICAL IMAGE', level=1)
            img_stream = io.BytesIO(image_bytes)
            doc.add_picture(img_stream, width=Inches(4))
        except:
            doc.add_paragraph("(Could not embed image)")
    
    # Analysis
    doc.add_heading('MEDICAL ANALYSIS', level=1)
    for para in analysis.split('\n'):
        if para.strip():
            doc.add_paragraph(para)
    
    # Disclaimer
    doc.add_paragraph().add_run(
        "This report is generated by artificial intelligence and should be reviewed by a qualified healthcare professional. It does not constitute professional medical advice."
    ).italic = True
    
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)
    return buffer

def generate_pptx(analysis, patient_info, image_bytes=None):
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = "MEDICAL ANALYSIS REPORT"
    date_str = datetime.now().strftime("%Y-%m-%d")
    subtitle.text = f"Generated on {date_str}"
    
    # Patient information slide
    bullet_slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Patient Information"
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    
    p = tf.paragraphs[0]
    p.text = f"Patient Name: {patient_info['name']}"
    
    p = tf.add_paragraph()
    p.text = f"Patient ID: {patient_info['id']}"
    
    p = tf.add_paragraph()
    p.text = f"Age: {patient_info['age']}"
    
    p = tf.add_paragraph()
    p.text = f"Sex: {patient_info['sex']}"
    
    # Image slide if available
    if image_bytes:
        try:
            img_slide_layout = prs.slide_layouts[6]  # Blank slide layout
            slide = prs.slides.add_slide(img_slide_layout)
            
            # Add title
            left = top = Inches(1)
            width = Inches(8)
            height = Inches(1)
            txBox = slide.shapes.add_textbox(left, top, width, height)
            tf = txBox.text_frame
            p = tf.add_paragraph()
            p.text = "Medical Image"
            p.font.size = Pt(24)
            p.font.bold = True
            
            # Add image
            img_stream = io.BytesIO(image_bytes)
            left = Inches(2)
            top = Inches(2)
            slide.shapes.add_picture(img_stream, left, top, height=Inches(4))
        except:
            pass
    
    # Analysis slides
    analysis_paragraphs = [p for p in analysis.split('\n') if p.strip()]
    chunks = [analysis_paragraphs[i:i+5] for i in range(0, len(analysis_paragraphs), 5)]
    
    for i, chunk in enumerate(chunks):
        if i == 0:
            slide_title = "Medical Analysis"
        else:
            slide_title = "Medical Analysis (Continued)"
            
        slide = prs.slides.add_slide(bullet_slide_layout)
        title = slide.shapes.title
        title.text = slide_title
        shapes = slide.shapes
        body_shape = shapes.placeholders[1]
        tf = body_shape.text_frame
        
        for para in chunk:
            p = tf.add_paragraph()
            p.text = para
            p.level = 0
    
    # Disclaimer slide
    slide = prs.slides.add_slide(bullet_slide_layout)
    title = slide.shapes.title
    title.text = "Disclaimer"
    shapes = slide.shapes
    body_shape = shapes.placeholders[1]
    tf = body_shape.text_frame
    p = tf.paragraphs[0]
    p.text = "This report is generated by artificial intelligence and should be reviewed by a qualified healthcare professional. It does not constitute professional medical advice."
    
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    return buffer

# Authentication UI
if not st.session_state.authenticated:
    col1, col2, col3 = st.columns([1,2,1])
    with col2:
        st.markdown("""
            <div style='text-align: center; padding: 20px;'>
                <h1 style='color: #1e3a8a; font-size: 2.5rem; margin-bottom: 1rem;'>Medical Report Generator</h1>
                <p style='font-size: 16px; color: #4b5563; margin-bottom: 30px;'>
                    Professional AI-powered medical imaging analysis
                </p>
            </div>
        """, unsafe_allow_html=True)
    
    st.markdown("<div class='auth-container'>", unsafe_allow_html=True)
    
    tab1, tab2 = st.tabs(["Login", "Sign Up"])
    
    with tab1:
        st.markdown("<h3 style='text-align: center; margin-bottom: 20px; color: #1e3a8a;'>Welcome Back</h3>", unsafe_allow_html=True)
        
        # Add Google Sign-In button using Streamlit's native button
        col1, col2, col3 = st.columns([1,3,1])
        with col2:
            # Add Google Sign-In container and script in the head
            st.components.v1.html("""
                <head>
                    <script src="https://accounts.google.com/gsi/client" async defer></script>
                </head>
                <body>
                    <div style="display: flex; justify-content: center; margin: 20px 0;">
                        <div id="g_id_onload"
                            data-client_id="1093379467339-kg7o7hbrnhn2rjmpr294k88fo5h7dpmh.apps.googleusercontent.com"
                            data-callback="handleCredentialResponse"
                            data-auto_prompt="false"
                            data-itp_support="true"
                            data-ux_mode="popup"
                            data-scope="openid email profile">
                        </div>
                        <div class="g_id_signin"
                            data-type="standard"
                            data-size="large"
                            data-theme="filled_black"
                            data-text="signin_with"
                            data-shape="rectangular"
                            data-logo_alignment="left"
                            data-width="300">
                        </div>
                    </div>
                    <script>
                        console.log("Current origin:", window.location.origin);
                        
                        function handleCredentialResponse(response) {
                            console.log("Received credential response");
                            if (response.credential) {
                                console.log("Got credential token");
                                
                                // Create form for submission
                                var form = document.createElement('form');
                                form.method = 'POST';
                                form.action = window.location.href;
                                
                                var input = document.createElement('input');
                                input.type = 'hidden';
                                input.name = 'google_token';
                                input.value = response.credential;
                                
                                form.appendChild(input);
                                document.body.appendChild(form);
                                console.log("Submitting form...");
                                form.submit();
                            } else {
                                console.error("No credential received in response:", response);
                            }
                        }
                    </script>
                </body>
            """, height=70)
        
        st.markdown("<div style='text-align: center; margin: 10px 0;'>or</div>", unsafe_allow_html=True)
        
        with st.form("login_form"):
            username = st.text_input("Email", placeholder="Enter your email")
            password = st.text_input("Password", type="password", placeholder="Enter your password")
            st.markdown("<br>", unsafe_allow_html=True)
            submit = st.form_submit_button("Login")
            
            if submit:
                if login(username, password):
                    st.markdown("""
                        <div class='success-message'>
                            <span class='medical-icon'>✓</span> Login successful! Redirecting...
                        </div>
                    """, unsafe_allow_html=True)
                    st.rerun()
                else:
                    st.markdown("""
                        <div class='error-message'>
                            <span class='medical-icon'>⚠️</span> Invalid username or password
                        </div>
                    """, unsafe_allow_html=True)
    
    with tab2:
        st.markdown("<h3 style='text-align: center; margin-bottom: 20px;'>Create New Account</h3>", unsafe_allow_html=True)
        
        with st.form("signup_form"):
            new_username = st.text_input("Username", placeholder="Choose a username")
            new_password = st.text_input("Password", type="password", placeholder="Choose a password")
            confirm_password = st.text_input("Confirm Password", type="password", placeholder="Confirm your password")
            email = st.text_input("Email", placeholder="Enter your email")
            st.markdown("<br>", unsafe_allow_html=True)
            submit_signup = st.form_submit_button("Sign Up")
            
            if submit_signup:
                if new_password != confirm_password:
                    st.markdown("""
                        <div class='error-message'>
                            Passwords do not match!
                        </div>
                    """, unsafe_allow_html=True)
                elif len(new_password) < 6:
                    st.markdown("""
                        <div class='error-message'>
                            Password must be at least 6 characters long!
                        </div>
                    """, unsafe_allow_html=True)
                else:
                    if create_user(new_username, new_password, email):
                        # Automatically log in the user after successful signup
                        if login(email, new_password):
                            st.markdown("""
                                <div class='success-message'>
                                    Account created successfully! Redirecting...
                                </div>
                            """, unsafe_allow_html=True)
                            st.rerun()
                        else:
                            st.markdown("""
                                <div class='error-message'>
                                    Account created but login failed. Please try logging in manually.
                                </div>
                            """, unsafe_allow_html=True)
                    else:
                        st.markdown("""
                            <div class='error-message'>
                                Username or email already exists!
                            </div>
                        """, unsafe_allow_html=True)
    
    st.markdown("</div>", unsafe_allow_html=True)
    
    # Add some testimonials or features section
    st.markdown("""
        <div style='max-width: 800px; margin: 3rem auto; text-align: center;'>
            <h3>Why Choose Our Medical Report Generator?</h3>
            <div style='display: flex; justify-content: space-between; margin-top: 30px;'>
                <div style='flex: 1; padding: 20px; background-color: white; border-radius: 8px; margin: 0 10px; box-shadow: 0 5px 15px rgba(0,0,0,0.05);'>
                    <h4>🔍 Advanced AI Analysis</h4>
                    <p>Powered by Google's Gemini AI for accurate medical image interpretation</p>
                </div>
                <div style='flex: 1; padding: 20px; background-color: white; border-radius: 8px; margin: 0 10px; box-shadow: 0 5px 15px rgba(0,0,0,0.05);'>
                    <h4>🌍 Multilingual Support</h4>
                    <p>Generate reports in 13 different languages for global accessibility</p>
                </div>
                <div style='flex: 1; padding: 20px; background-color: white; border-radius: 8px; margin: 0 10px; box-shadow: 0 5px 15px rgba(0,0,0,0.05);'>
                    <h4>🔒 Secure & Private</h4>
                    <p>Your medical data is processed with the highest security standards</p>
                </div>
            </div>
        </div>
    """, unsafe_allow_html=True)
    
    st.stop()

# Main application (only shown after authentication)
st.markdown("""
    <div class='report-header'>
        <h1 style='margin-bottom: 0.5rem;'>Medical Report Generator</h1>
        <p style='color: #4b5563; margin-bottom: 0;'>Professional medical imaging analysis and reporting</p>
    </div>
""", unsafe_allow_html=True)

# User profile section
st.markdown(f"""
    <div class='profile-section'>
        <div>
            <h3 style='margin-bottom: 5px;'>Welcome, {st.session_state.username}!</h3>
            <p style='color: #6c757d;'>Let's generate some professional medical reports today.</p>
        </div>
        <div>
            <button class='stButton' style='background: #f8f9fa; color: #5B86E5; border: 1px solid #5B86E5; box-shadow: none; padding: 8px 16px; font-size: 14px; border-radius: 4px;' onclick="logout()">Logout</button>
        </div>
    </div>
""", unsafe_allow_html=True)

# Two column layout for better organization
col1, col2 = st.columns([1, 2])

with col1:
    st.markdown("<div class='language-selector'>", unsafe_allow_html=True)
    st.markdown("<h4>Report Settings</h4>", unsafe_allow_html=True)
    
    # Language selection
    selected_language = st.selectbox(
        "Select Report Language",
        options=list(LANGUAGES.keys()),
        index=list(LANGUAGES.keys()).index(st.session_state.language)
    )
    st.session_state.language = selected_language
    
    # Remove API Key input and use the provided key
    API_KEY = "AIzaSyBjVRQNCTl_HykJb0WCqvO_T7JI7FTkRmE"
    
    # Add logout button
    if st.button("Logout"):
        logout()
        st.rerun()
    
    st.markdown("</div>", unsafe_allow_html=True)
    
    # Patient information form
    st.markdown("<div class='patient-info-form'>", unsafe_allow_html=True)
    st.markdown("<h4 class='section-header'>Patient Information</h4>", unsafe_allow_html=True)
    
    patient_name = st.text_input("Patient Name", placeholder="Enter patient's full name")
    patient_id = st.text_input("Patient ID", placeholder="Enter patient's ID number")
    patient_age = st.number_input("Age", min_value=0, max_value=120, step=1)
    patient_sex = st.selectbox("Sex", ["Male", "Female", "Other"])
    
    # Store patient info in session state
    st.session_state.patient_info = {
        'name': patient_name,
        'id': patient_id,
        'age': patient_age,
        'sex': patient_sex
    }
    
    st.markdown("</div>", unsafe_allow_html=True)
    
    # Information box
    st.markdown("""
        <div style='background-color: #e3f2fd; padding: 20px; border-radius: 8px; margin-top: 20px;'>
            <h4 style='color: #0d47a1;'>How It Works</h4>
            <ol style='padding-left: 20px;'>
                <li>Upload a medical image (X-ray, MRI, CT scan, etc.)</li>
                <li>Enter patient information</li>
                <li>Click "Generate Report"</li>
                <li>Receive a detailed medical analysis</li>
                <li>Download the report in your preferred format</li>
            </ol>
            <p style='font-style: italic; margin-top: 15px; font-size: 14px;'>This tool uses Google's Gemini AI to analyze medical images and provide detailed medical descriptions.</p>
        </div>
    """, unsafe_allow_html=True)

with col2:
    # File uploader with better styling
    st.markdown("<h3>Upload Medical Image</h3>", unsafe_allow_html=True)
    
    uploaded_file = st.file_uploader("Choose a medical image...", type=['jpg', 'jpeg', 'png'])
    
    if uploaded_file is not None:
        # Display the uploaded image in a nice frame
        st.markdown("<div class='uploadedFile'>", unsafe_allow_html=True)
        image = Image.open(uploaded_file)
        st.image(image, caption="Uploaded Medical Image", use_column_width=True)
        st.markdown("</div>", unsafe_allow_html=True)
        
        # Store image bytes in session state
        st.session_state.image_bytes = uploaded_file.getvalue()
        
        # Analysis button
        if st.session_state.patient_info['name'] and st.session_state.patient_info['id']:
            if st.button("Generate Report", key="generate_btn"):
                with st.spinner("Analyzing image..."):
                    try:
                        # Convert image to base64
                        image_base64 = base64.b64encode(st.session_state.image_bytes).decode()
                        
                        # Prepare the request
                        url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={API_KEY}"
                        
                        headers = {
                            "Content-Type": "application/json"
                        }
                        
                        # Get the appropriate prompt for the selected language
                        language_code = LANGUAGES[selected_language]
                        prompt = LANGUAGE_PROMPTS[language_code]
                        
                        data = {
                            "contents": [{
                                "parts": [
                                    {
                                        "inline_data": {
                                            "mime_type": uploaded_file.type,
                                            "data": image_base64
                                        }
                                    },
                                    {
                                        "text": prompt
                                    }
                                ]
                            }]
                        }
                        
                        # Make the API request
                        response = requests.post(url, headers=headers, json=data)
                        response.raise_for_status()
                        
                        # Extract and display the analysis
                        result = response.json()
                        analysis = result['candidates'][0]['content']['parts'][0]['text']
                        
                        # Store analysis in session state
                        st.session_state.analysis_result = analysis
                        
                        # Display the report in a nice format
                        st.markdown("<div class='report-card'>", unsafe_allow_html=True)
                        st.markdown(f"<h3 style='color: #1e3a8a;'>📋 Medical Analysis Report ({selected_language})</h3>", unsafe_allow_html=True)
                        st.markdown("<hr style='border-color: #e2e8f0; margin: 1rem 0;'>", unsafe_allow_html=True)
                        st.markdown(analysis)
                        st.markdown("</div>", unsafe_allow_html=True)
                        
                    except Exception as e:
                        st.markdown(f"""
                            <div class='error-message'>
                                <h4>⚠️ Error</h4>
                                <p>{str(e)}</p>
                                <p>Please check your API key and try again.</p>
                            </div>
                        """, unsafe_allow_html=True)
        elif uploaded_file is not None:
            if not API_KEY:
                st.warning("Please enter your Google AI Studio API key to proceed with the analysis.")
            if not st.session_state.patient_info['name'] or not st.session_state.patient_info['id']:
                st.warning("Please complete the patient information form to proceed with the analysis.")
    else:
        # Display a placeholder or instructions
        st.markdown("""
            <div style='background-color: #f8f9fa; border: 2px dashed #dddddd; border-radius: 10px; padding: 50px 20px; text-align: center; margin-top: 20px;'>
                <img src="https://via.placeholder.com/150/5B86E5/FFFFFF?text=📊" style='width: 100px; height: 100px; opacity: 0.7;'>
                <h3 style='margin-top: 20px; color: #6c757d;'>Upload a Medical Image</h3>
                <p style='color: #6c757d;'>Supported formats: JPG, JPEG, PNG</p>
            </div>
        """, unsafe_allow_html=True)

    # Download options (only shown if there's an analysis result)
    if st.session_state.analysis_result and st.session_state.patient_info and st.session_state.image_bytes:
        st.markdown("<h4>Download Report</h4>", unsafe_allow_html=True)
        st.markdown("<div class='download-options'>", unsafe_allow_html=True)
        
        col1, col2, col3 = st.columns(3)
        
        with col1:
            # PDF download
            pdf_buffer = generate_pdf(
                st.session_state.analysis_result,
                st.session_state.patient_info,
                st.session_state.image_bytes
            )
            st.download_button(
                label="📄 Download PDF",
                data=pdf_buffer,
                file_name="medical_report.pdf",
                mime="application/pdf",
                key="pdf_download"
            )
        
        with col2:
            # Word download
            docx_buffer = generate_docx(
                st.session_state.analysis_result,
                st.session_state.patient_info,
                st.session_state.image_bytes
            )
            st.download_button(
                label="📝 Download Word",
                data=docx_buffer,
                file_name="medical_report.docx",
                mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                key="docx_download"
            )
        
        with col3:
            # PowerPoint download
            pptx_buffer = generate_pptx(
                st.session_state.analysis_result,
                st.session_state.patient_info,
                st.session_state.image_bytes
            )
            st.download_button(
                label="📊 Download PowerPoint",
                data=pptx_buffer,
                file_name="medical_report.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                key="pptx_download"
            )
        
        # Also provide the plain text option
        st.download_button(
            label="📄 Download Plain Text",
            data=st.session_state.analysis_result,
            file_name="medical_report.txt",
            mime="text/plain",
            key="txt_download"
        )
        
        st.markdown("</div>", unsafe_allow_html=True)

# Footer with better styling
st.markdown("""
    <div class='footer'>
        <div style='display: flex; justify-content: center; margin-bottom: 15px;'>
            <div style='margin: 0 15px; color: #4b5563;'>Terms of Service</div>
            <div style='margin: 0 15px; color: #4b5563;'>Privacy Policy</div>
            <div style='margin: 0 15px; color: #4b5563;'>FAQs</div>
            <div style='margin: 0 15px; color: #4b5563;'>Contact</div>
        </div>
        <p style='color: #4b5563;'>© 2025 Medical Report Generator. Professional medical imaging analysis platform</p>
        <p style='font-size: 12px; color: #6b7280; margin-top: 10px;'>This tool is intended for informational purposes only and should not replace professional medical advice.</p>
    </div>
""", unsafe_allow_html=True)

# Handle Google Sign-In response
if 'google_token' in st.session_state or st.experimental_get_query_params().get('google_token'):
    st.write("Processing Google Sign-In...")
    token = st.session_state.get('google_token') or st.experimental_get_query_params().get('google_token')[0]
    if handle_google_signin(token):
        st.success("Successfully signed in with Google!")
        time.sleep(1)
        st.rerun()
    else:
        st.error("Failed to process Google Sign-In. Please try again.")